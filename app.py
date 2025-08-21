import os
import io
import requests
import html2text
from urllib.parse import urlparse
from flask import Flask, render_template, request, Response, send_file
from lxml import etree, html as lxml_html
from docx import Document
from docx.shared import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation

app = Flask(__name__)

XSL_PATH = os.path.join(os.path.dirname(__file__), "xsl", "rss_substack_fragment.xsl")

def parse_posts_from_cards(html_fragment: str):
    try:
        doc = lxml_html.fromstring(html_fragment)
        posts = []
        for card in doc.xpath('//div[@class="card"]'):
            title_element = card.find('.//h2/a')
            content_element = card.find('.//div[@class="content"]')
            title = title_element.text_content().strip() if title_element is not None else "Untitled Post"
            if content_element is not None:
                etree.strip_elements(content_element, "script", "style", with_tail=False)
                content = "\n\n".join([line.strip() for line in content_element.text_content().splitlines() if line.strip()])
            else: content = "No content available."
            posts.append({'title': title, 'content': content})
        return posts
    except Exception:
        return [{'title': 'Feed Content', 'content': html_to_text(html_fragment)}]

def normalize_substack_input(text: str):
    text = (text or "").strip()
    if not text: return ""
    if "substack.com" not in text: text = f"{text}.substack.com"
    if not text.startswith("http://") and not text.startswith("https://"): text = "https://" + text
    parsed = urlparse(text)
    path = parsed.path.rstrip("/")
    if not path.endswith("/feed"): path = path + "/feed"
    return f"{parsed.scheme}://{parsed.netloc}{path}"

def fetch_xml(url: str):
    headers = {"User-Agent": "SubstackRSSViewer/2.0", "Accept": "application/rss+xml, application/xml"}
    r = requests.get(url, headers=headers, timeout=25)
    r.raise_for_status()
    return r.content

def transform_to_fragment(xml_content: bytes):
    with open(XSL_PATH, "rb") as f: xsl_doc = etree.XML(f.read())
    transform = etree.XSLT(xsl_doc)
    parser = etree.XMLParser(recover=True, resolve_entities=False)
    xml_doc = etree.fromstring(xml_content, parser=parser)
    return str(transform(xml_doc))

def xml_channel_title(xml_content: bytes):
    try:
        doc = etree.fromstring(xml_content)
        t = doc.findtext(".//channel/title")
        return t.strip() if t else "Feed"
    except Exception:
        return "Feed"

def html_to_text(html_str: str):
    try:
        doc = lxml_html.fromstring(html_str)
        etree.strip_elements(doc, "script", "style", with_tail=False)
        return "\n".join([line.strip() for line in doc.text_content().splitlines() if line.strip()])
    except Exception:
        return ""

@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")

@app.route("/docs")
def docs():
    return render_template("docs.html")

@app.route("/view", methods=["POST"])
def view_feed():
    raw = request.form.get("feed_input", "").strip()
    normalized = normalize_substack_input(raw)
    if not normalized:
        return render_template("index.html", error="Please enter a valid Substack name or URL.", last=raw)
    try:
        xml = fetch_xml(normalized)
        fragment = transform_to_fragment(xml)
        title = xml_channel_title(xml)
        return render_template("result.html", feed_html=fragment, feed_url=normalized, feed_title=title)
    except Exception as e:
        return render_template("index.html", error=f"Couldn't fetch the feed: {e}", last=raw)

class FeedDownloader:
    def __init__(self, feed_url):
        self.feed_url = feed_url
        if not self.feed_url: raise ValueError("Missing url parameter")
        self.xml = fetch_xml(self.feed_url)
        self.channel_title = xml_channel_title(self.xml)
        self.html_fragment = transform_to_fragment(self.xml)
        self.posts = parse_posts_from_cards(self.html_fragment)
        self.safe_title = "".join(c for c in self.channel_title if c.isalnum() or c in " -_").rstrip()

@app.route("/download/pdf")
def download_pdf():
    try:
        d = FeedDownloader(request.args.get("url", ""))
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(name='PostTitle', parent=styles['h2'], spaceAfter=14))
        story = [Paragraph(d.channel_title, styles['h1']), Spacer(1, 0.25 * inch)]
        for post in d.posts:
            story.append(Paragraph(post['title'], styles['PostTitle']))
            for p_text in post['content'].split('\n\n'):
                story.append(Paragraph(p_text, styles['BodyText']))
                story.append(Spacer(1, 0.1 * inch))
            story.append(Spacer(1, 0.5 * inch))
        doc.build(story)
        buffer.seek(0)
        return send_file(buffer, as_attachment=True, download_name=f"{d.safe_title}.pdf", mimetype="application/pdf")
    except Exception as e: return Response(f"Error: {e}", status=500)

@app.route("/download/docx")
def download_docx():
    try:
        d = FeedDownloader(request.args.get("url", ""))
        doc = Document()
        for section in doc.sections: section.left_margin, section.right_margin = Inches(1), Inches(1)
        doc.add_heading(d.channel_title, level=1)
        for post in d.posts:
            doc.add_heading(post['title'], level=2)
            for para in post['content'].split('\n\n'): doc.add_paragraph(para, style='Body Text')
            p = doc.add_paragraph(); p.paragraph_format.space_before = Pt(12)
        buffer = io.BytesIO(); doc.save(buffer); buffer.seek(0)
        return send_file(buffer, as_attachment=True, download_name=f"{d.safe_title}.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    except Exception as e: return Response(f"Error: {e}", status=500)

@app.route("/download/pptx")
def download_pptx():
    try:
        d = FeedDownloader(request.args.get("url", ""))
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text = d.channel_title
        for post in d.posts:
            slide = prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text = post['title']
            slide.placeholders[1].text = post['content']
        buffer = io.BytesIO(); prs.save(buffer); buffer.seek(0)
        return send_file(buffer, as_attachment=True, download_name=f"{d.safe_title}.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e: return Response(f"Error: {e}", status=500)

@app.route("/download/txt")
def download_txt():
    try:
        d = FeedDownloader(request.args.get("url", ""))
        full_text = f"{d.channel_title}\n{'='*len(d.channel_title)}\n\n"
        for post in d.posts: full_text += f"{post['title']}\n{'-'*len(post['title'])}\n\n{post['content']}\n\n\n"
        return send_file(io.BytesIO(full_text.encode('utf-8')), as_attachment=True, download_name=f"{d.safe_title}.txt", mimetype="text/plain")
    except Exception as e: return Response(f"Error: {e}", status=500)

@app.route("/download/md")
def download_md():
    try:
        d = FeedDownloader(request.args.get("url", ""))
        h = html2text.HTML2Text(); h.body_width = 0
        markdown = h.handle(d.html_fragment)
        return send_file(io.BytesIO(markdown.encode('utf-8')), as_attachment=True, download_name=f"{d.safe_title}.md", mimetype="text/markdown")
    except Exception as e: return Response(f"Error: {e}", status=500)

@app.route("/download/html")
def download_html():
    try:
        d = FeedDownloader(request.args.get("url", ""))
        full_html = f'<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8"><title>{d.channel_title}</title><style>body{{font-family:sans-serif;max-width:800px;margin:2rem auto;}}</style></head><body>{d.html_fragment}</body></html>'
        return send_file(io.BytesIO(full_html.encode('utf-8')), as_attachment=True, download_name=f"{d.safe_title}.html", mimetype="text/html")
    except Exception as e: return Response(f"Error: {e}", status=500)

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
